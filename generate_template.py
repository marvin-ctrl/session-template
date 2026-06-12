#!/usr/bin/env python3
"""
Futsal National Team — Session Plan Template Generator
======================================================
Generates `futsal_session_template.pptx`: a fully editable, premium,
one-page session-plan template system for national-team futsal staff.

Slides:
  1. Master template       (keep clean, duplicate weekly)
  2. Working copy          (identical, ready to fill)
  3. Annotated guide       (same layout + deletable "HOW TO" callouts)
  4. Component library     (reusable futsal tactical symbols + courts)
  5. Session control board (time, load, tactical balance, rotations,
                            review loop — the session as an instrument)
  6. On-court cards        (print, cut, pocket — arm's-length cue cards)

Everything is native PowerPoint vector shapes / text boxes / tables.
No images, no locked elements, no hidden masking shapes.

Usage:  python3 generate_template.py
"""

import copy

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- constants

EMU_IN = 914400


def IN(v):
    """Inches -> Emu."""
    return Emu(int(round(v * EMU_IN)))


SLIDE_W = 13.333
SLIDE_H = 7.5

C_DARK = RGBColor(0x1A, 0x1A, 0x1A)    # headings, primary text
C_MID = RGBColor(0x6B, 0x6B, 0x6B)     # secondary text
C_LITE = RGBColor(0xF2, 0xF2, 0xF2)    # focus bar / panel fills
C_FAINT = RGBColor(0xFA, 0xFA, 0xFA)   # subtle row tint
C_LINE = RGBColor(0xD9, 0xD9, 0xD9)    # hairline dividers
C_ACCENT = RGBColor(0x00, 0x57, 0xB8)  # single accent (hierarchy only)
C_COURT = RGBColor(0xA9, 0xB1, 0xBB)   # court line grey
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

# Court-scale marker size.
# On a 3.26" wide court (=40m), 1m = 0.0815".
# 0.14" ≈ 1.7 m — standard for elite coaching diagrams.
MARKER_D = 0.14

# page geometry (all inches; arithmetic validated to sum exactly)
MARGIN = 0.15
HEADER_H = 0.55
RULE_H = 0.02
FOCUS_Y = HEADER_H + RULE_H            # 0.57
FOCUS_H = 0.65
BODY_Y = 1.28
BODY_H = 6.07                          # bottom edge 7.35
BLOCK_W = 5.265
BLOCK_H = 3.00
BLOCK_GAP = 0.07
SQUAD_X = 10.83
SQUAD_W = 2.35

_shape_id = [9000]                     # id counter for hand-built group XML


def _next_id():
    _shape_id[0] += 1
    return _shape_id[0]


# ---------------------------------------------------------------- low-level helpers

def _no_shadow(shp):
    shp.shadow.inherit = False
    return shp


def add_box(shapes, x, y, w, h, fill=None, line=None, line_w=0.75,
            shape=MSO_SHAPE.RECTANGLE, dash=None):
    """Add a rectangle/autoshape with explicit fill + outline (None = off)."""
    shp = shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
        if dash is not None:
            shp.line.dash_style = dash
    return _no_shadow(shp)


def add_line(shapes, x1, y1, x2, y2, color=C_LINE, w=0.75, dash=None):
    ln = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                              IN(x1), IN(y1), IN(x2), IN(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    if dash is not None:
        ln.line.dash_style = dash
    return _no_shadow(ln)


def add_arrowhead(line_shape):
    """python-pptx has no arrowhead API; append <a:tailEnd> to the line XML."""
    ln = line_shape.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"),
                          {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return line_shape


def set_fill_alpha(shp, opacity_pct):
    """Set transparency on a solid fill via <a:alpha> (no python-pptx API)."""
    solid = shp.fill._xPr.find(qn("a:solidFill"))
    clr = solid.find(qn("a:srgbClr"))
    clr.append(clr.makeelement(qn("a:alpha"),
                               {"val": str(int(opacity_pct * 1000))}))


def add_text(shapes, x, y, w, h, paras, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, wrap=True):
    """
    Text box. `paras` = list of paragraphs; each paragraph = list of run
    tuples (text, size_pt, bold, color). One-click editable, no autofit.
    """
    tb = shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = IN(0.03)
    tf.margin_right = IN(0.03)
    tf.margin_top = IN(0.015)
    tf.margin_bottom = IN(0.015)
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for text, size, bold, color in runs:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = FONT
            f.size = Pt(size)
            f.bold = bold
            f.color.rgb = color
            f.language_id = MSO_LANGUAGE_ID.ENGLISH_UK
    return _no_shadow(tb)


def group_elements(slide, shps, x, y, w, h, name):
    """Wrap already-placed shapes into a <p:grpSp> with an identity
    child-coordinate transform, so the group moves/copies as one unit."""
    ex, ey, ew, eh = (IN(x), IN(y), IN(w), IN(h))
    grp = parse_xml(
        '<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvGrpSpPr><p:cNvPr id="{_next_id()}" name="{name}"/>'
        '<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        '<p:grpSpPr><a:xfrm>'
        f'<a:off x="{ex}" y="{ey}"/><a:ext cx="{ew}" cy="{eh}"/>'
        f'<a:chOff x="{ex}" y="{ey}"/><a:chExt cx="{ew}" cy="{eh}"/>'
        '</a:xfrm></p:grpSpPr></p:grpSp>'
    )
    slide.shapes._spTree.append(grp)
    for s in shps:
        grp.append(s._element)
    return grp


def set_adj_angles(shp, start_deg, end_deg):
    """Set adj1/adj2 angle adjustments (clockwise from 3 o'clock) via avLst
    XML — python-pptx exposes no adjustment spec for arc/pie/chord presets."""
    prst = shp._element.spPr.find(qn("a:prstGeom"))
    av = prst.find(qn("a:avLst"))
    if av is None:
        av = prst.makeelement(qn("a:avLst"), {})
        prst.append(av)
    else:
        for gd in list(av):
            av.remove(gd)
    for nm, deg in (("adj1", start_deg), ("adj2", end_deg)):
        gd = av.makeelement(qn("a:gd"), {"name": nm,
                                         "fmla": f"val {int(deg * 60000)}"})
        av.append(gd)
    return shp


def add_arc(shapes, cx, cy, r, start_deg, end_deg, color=C_COURT, w=0.75):
    """MSO ARC autoshape with explicit start/end angles."""
    shp = shapes.add_shape(MSO_SHAPE.ARC,
                           IN(cx - r), IN(cy - r), IN(2 * r), IN(2 * r))
    set_adj_angles(shp, start_deg, end_deg)
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(w)
    return _no_shadow(shp)


# ---------------------------------------------------------------- futsal court

def draw_futsal_court(slide, x, y, w, group_name="Futsal court"):
    """
    Futsal-correct court at true 40m x 20m proportions (h = w/2).
    Penalty areas: two 6m quarter-arcs centred on each goalpost joined by a
    straight 3m segment (the real futsal D — not a football box, not a
    plain semicircle). Returns the grouped court.
    """
    h = w / 2.0
    s = w / 40.0                       # inches per metre
    cy = y + h / 2.0
    cx = x + w / 2.0
    goal_d = 0.07                      # goal stub depth behind goal line
    dot = 0.035                        # mark diameter
    parts = []

    # touchlines / goal lines
    parts.append(add_box(slide.shapes, x, y, w, h, fill=None,
                         line=C_COURT, line_w=1.0))
    # halfway line + centre circle + centre mark
    parts.append(add_line(slide.shapes, cx, y, cx, y + h, C_COURT, 0.75))
    parts.append(add_box(slide.shapes, cx - 3 * s, cy - 3 * s, 6 * s, 6 * s,
                         fill=None, line=C_COURT, line_w=0.75,
                         shape=MSO_SHAPE.OVAL))
    parts.append(add_box(slide.shapes, cx - dot / 2, cy - dot / 2, dot, dot,
                         fill=C_COURT, line=None, shape=MSO_SHAPE.OVAL))

    for end in ("L", "R"):
        gl = x if end == "L" else x + w          # goal line x
        sgn = 1 if end == "L" else -1            # into-court direction
        post_top = y + 8.5 * s                   # goal is 3m, centred
        post_bot = y + 11.5 * s
        # penalty area: quarter arc off each post + straight 6m segment
        if end == "L":
            parts.append(add_arc(slide.shapes, gl, post_top, 6 * s, 270, 360))
            parts.append(add_arc(slide.shapes, gl, post_bot, 6 * s, 0, 90))
        else:
            parts.append(add_arc(slide.shapes, gl, post_top, 6 * s, 180, 270))
            parts.append(add_arc(slide.shapes, gl, post_bot, 6 * s, 90, 180))
        seg_x = gl + sgn * 6 * s
        parts.append(add_line(slide.shapes, seg_x, post_top, seg_x, post_bot,
                              C_COURT, 0.75))
        # 6m penalty mark + 10m second penalty mark
        for dist in (6, 10):
            mx = gl + sgn * dist * s
            parts.append(add_box(slide.shapes, mx - dot / 2, cy - dot / 2,
                                 dot, dot, fill=C_COURT, line=None,
                                 shape=MSO_SHAPE.OVAL))
        # goal (3m wide stub behind the goal line)
        gx = gl - goal_d if end == "L" else gl
        parts.append(add_box(slide.shapes, gx, post_top, goal_d, 3 * s,
                             fill=None, line=C_COURT, line_w=1.0))

    return group_elements(slide, parts, x - goal_d, y,
                          w + 2 * goal_d, h, group_name)


def draw_futsal_half_court(slide, x, y, w, group_name="Futsal half court"):
    """Half court (20m x 20m square): goal end left, halfway line right."""
    s = w / 20.0
    h = w
    cy = y + h / 2.0
    dot = 0.035
    goal_d = 0.07
    parts = []

    parts.append(add_box(slide.shapes, x, y, w, h, fill=None,
                         line=C_COURT, line_w=1.0))
    # half centre circle on the halfway (right) edge
    parts.append(add_arc(slide.shapes, x + w, cy, 3 * s, 90, 270))
    parts.append(add_box(slide.shapes, x + w - dot / 2, cy - dot / 2,
                         dot, dot, fill=C_COURT, line=None,
                         shape=MSO_SHAPE.OVAL))
    post_top = y + h / 2 - 1.5 * s
    post_bot = y + h / 2 + 1.5 * s
    parts.append(add_arc(slide.shapes, x, post_top, 6 * s, 270, 360))
    parts.append(add_arc(slide.shapes, x, post_bot, 6 * s, 0, 90))
    parts.append(add_line(slide.shapes, x + 6 * s, post_top,
                          x + 6 * s, post_bot, C_COURT, 0.75))
    for dist in (6, 10):
        parts.append(add_box(slide.shapes, x + dist * s - dot / 2,
                             cy - dot / 2, dot, dot,
                             fill=C_COURT, line=None, shape=MSO_SHAPE.OVAL))
    parts.append(add_box(slide.shapes, x - goal_d, post_top, goal_d, 3 * s,
                         fill=None, line=C_COURT, line_w=1.0))

    return group_elements(slide, parts, x - goal_d, y, w + goal_d, h,
                          group_name)


# ---------------------------------------------------------------- page zones

def draw_header(slide):
    """White masthead: team + session title left, logistics fields right."""
    sh = slide.shapes
    add_text(sh, MARGIN, 0.02, 4.30, 0.51, [
        [("[INSERT TEAM]  ·  FUTSAL SESSION PLAN", 8, True, C_MID)],
        [("[Insert session title]", 20, True, C_DARK)],
    ])
    fields = [
        ("DATE", "[Insert date]", 4.62, 1.45),
        ("VENUE", "[Insert venue]", 6.17, 1.95),
        ("DURATION", "[Insert duration]", 8.22, 1.45),
        ("STAFF", "[Insert staff]", 9.77, 3.41),
    ]
    for label, value, fx, fw in fields:
        add_text(sh, fx, 0.07, fw, 0.44, [
            [(label, 7, True, C_MID)],
            [(value, 9.5, label == "DURATION", C_DARK)],
        ])
    # accent rule separates masthead from the working document
    add_box(sh, 0, HEADER_H, SLIDE_W, RULE_H, fill=C_ACCENT, line=None)


def draw_focus_bar(slide):
    """Tier-1 focus strip: theme, outcomes, constraints, reminders."""
    sh = slide.shapes
    add_box(sh, 0, FOCUS_Y, SLIDE_W, FOCUS_H, fill=C_LITE, line=None)
    cols = [
        ("SESSION THEME / OBJECTIVE", "[Insert focus]", MARGIN, 3.30, True),
        ("KEY COACHING OUTCOMES", "[Insert key outcomes]", 3.60, 3.45, False),
        ("KEY CONSTRAINTS", "[Insert constraints]", 7.20, 3.00, False),
        ("KEY REMINDERS", "[Insert reminders]", 10.35, 2.83, False),
    ]
    for label, value, fx, fw, is_theme in cols:
        add_text(sh, fx, FOCUS_Y + 0.05, fw, FOCUS_H - 0.10, [
            [(label, 7, True, C_ACCENT if is_theme else C_MID)],
            [(value, 11 if is_theme else 8.5, is_theme,
              C_DARK if is_theme else C_DARK)],
        ])
        if fx > MARGIN:
            add_line(sh, fx - 0.08, FOCUS_Y + 0.10,
                     fx - 0.08, FOCUS_Y + FOCUS_H - 0.10, C_LINE, 0.75)


def draw_activity_block(slide, x, y, number, title):
    """One session block: title bar / objective / court + coaching points /
    notes. Identical structure for all four blocks."""
    sh = slide.shapes
    w, h = BLOCK_W, BLOCK_H
    # container + title bar
    add_box(sh, x, y, w, h, fill=C_WHITE, line=C_LINE, line_w=0.75)
    bar_h = 0.28
    add_box(sh, x, y, w, bar_h, fill=C_LITE, line=None)
    add_box(sh, x, y, 0.05, bar_h, fill=C_ACCENT, line=None)
    add_text(sh, x + 0.07, y + 0.015, 0.36, 0.25,
             [[(number, 11, True, C_ACCENT)]])
    add_text(sh, x + 0.44, y + 0.03, 1.82, 0.23,
             [[(title, 9, True, C_DARK)]], wrap=False)
    add_text(sh, x + 2.30, y + 0.035, w - 1.40 - 2.30, 0.22, [[
        ("LEAD  ", 7, True, C_MID),
        ("[Insert responsibility]", 8, False, C_DARK),
    ]], align=PP_ALIGN.RIGHT, wrap=False)
    chip = add_box(sh, x + w - 1.30, y + 0.035, 1.22, 0.21,
                   fill=C_WHITE, line=C_ACCENT, line_w=1.0,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = chip.text_frame
    tf.word_wrap = False
    tf.margin_left = IN(0.02)
    tf.margin_right = IN(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "[Insert timing]"
    r.font.name = FONT
    r.font.size = Pt(8)
    r.font.bold = True
    r.font.color.rgb = C_ACCENT

    # objective row
    obj_y = y + bar_h + 0.02
    add_text(sh, x + 0.07, obj_y, w - 0.14, 0.20, [[
        ("OBJECTIVE  ", 7, True, C_MID),
        ("[Insert objective]", 8.5, False, C_DARK),
    ]])

    # main area: court left / coaching points right
    main_y = obj_y + 0.21
    main_h = 1.92
    court_w = 3.26
    court_h = court_w / 2
    draw_futsal_court(slide, x + 0.12, main_y + (main_h - court_h) / 2,
                      court_w, f"Court {number}")
    div_x = x + 0.12 + court_w + 0.12
    add_line(sh, div_x, main_y + 0.06, div_x, main_y + main_h - 0.06,
             C_LINE, 0.75)
    cp_x = div_x + 0.06
    cp_w = x + w - 0.07 - cp_x
    add_text(sh, cp_x, main_y + 0.03, cp_w, main_h - 0.06, [
        [("COACHING POINTS", 7, True, C_MID)],
        [("[Insert coaching points]", 8, False, C_DARK)],
    ])

    # notes strip
    notes_y = main_y + main_h + 0.02
    notes_h = y + h - notes_y - 0.04
    add_line(sh, x + 0.07, notes_y, x + w - 0.07, notes_y, C_LINE, 0.75)
    add_text(sh, x + 0.07, notes_y + 0.02, w - 0.14, notes_h, [[
        ("NOTES / CONSTRAINTS / PROGRESSIONS   ", 7, True, C_MID),
        ("[Insert notes]", 8, False, C_DARK),
    ]])


def _style_cell(cell, text, size, bold, color, fill, align=PP_ALIGN.LEFT):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = IN(0.04)
    cell.margin_right = IN(0.02)
    cell.margin_top = IN(0.005)
    cell.margin_bottom = IN(0.005)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color


def draw_squad_panel(slide):
    """Operational availability panel: counts, position-seeded player table
    (tab-through editing), squad notes."""
    sh = slide.shapes
    x, y, w, h = SQUAD_X, BODY_Y, SQUAD_W, BODY_H
    add_box(sh, x, y, w, h, fill=C_WHITE, line=C_LINE, line_w=0.75)
    add_box(sh, x, y, w, 0.30, fill=C_DARK, line=None)
    add_text(sh, x + 0.05, y + 0.045, w - 0.10, 0.24,
             [[("AVAILABLE PLAYERS / POSITIONS", 8, True, C_WHITE)]],
             wrap=False)
    # counts row — squad balance at a glance
    add_text(sh, x + 0.05, y + 0.33, w - 0.10, 0.24, [[
        ("TOTAL ", 7, True, C_MID), ("[#]   ", 9, True, C_DARK),
        ("GK ", 7, True, C_MID), ("[#]   ", 9, True, C_DARK),
        ("OUTFIELD ", 7, True, C_MID), ("[#]", 9, True, C_DARK),
    ]])

    # player table: header + 14 rows seeded with futsal positions
    positions = (["GK"] * 2 + ["FIXO"] * 3 + ["ALA"] * 6 + ["PIVOT"] * 3)
    rows = 1 + len(positions)
    tbl_y = y + 0.60
    tbl_h = 4.05                       # 15 rows x 0.27"; slack below is
    #                                    deliberate room for late additions
    gframe = sh.add_table(rows, 3, IN(x + 0.05), IN(tbl_y),
                          IN(w - 0.10), IN(tbl_h))
    tbl = gframe.table
    tbl.first_row = False
    tbl.horz_banding = False
    # "No Style, No Grid" — fills only, no theme-blue styling
    tblPr = tbl._tbl.tblPr
    for el in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(el)
    tsid = tblPr.makeelement(qn("a:tableStyleId"), {})
    tsid.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
    tblPr.append(tsid)

    tbl.columns[0].width = IN(0.52)
    tbl.columns[1].width = IN(1.08)
    tbl.columns[2].width = IN(0.65)
    row_h = IN(tbl_h / rows)
    for r in tbl.rows:
        r.height = row_h

    for c, head in enumerate(("POS", "PLAYER", "STATUS")):
        _style_cell(tbl.cell(0, c), head, 7, True, C_DARK, C_LITE)
    for i, pos in enumerate(positions, start=1):
        tint = C_FAINT if pos in ("FIXO", "PIVOT") else C_WHITE
        _style_cell(tbl.cell(i, 0), pos, 7.5, True, C_DARK, tint)
        _style_cell(tbl.cell(i, 1), "[Player]", 7.5, False, C_MID, tint)
        _style_cell(tbl.cell(i, 2), "", 7.5, False, C_MID, tint)

    # squad notes (anchored to panel bottom, clear of any table reflow)
    notes_y = y + 5.25
    add_line(sh, x + 0.05, notes_y, x + w - 0.05, notes_y, C_LINE, 0.75)
    add_text(sh, x + 0.05, notes_y + 0.03, w - 0.10,
             y + h - notes_y - 0.08, [
                 [("SQUAD NOTES", 7, True, C_MID)],
                 [("[Insert status notes]", 8, False, C_DARK)],
             ])


def add_footer(slide, text):
    add_text(slide.shapes, MARGIN, 7.355, SLIDE_W - 2 * MARGIN, 0.135,
             [[(text, 6.5, False, C_MID)]])


# ---------------------------------------------------------------- slide builders

BLOCKS = [
    ("01", "WARM UP"),
    ("02", "TECHNICAL / POSSESSION"),
    ("03", "TACTICAL ACTIVITY"),
    ("04", "GAME"),
]


def build_session_slide(prs, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
    draw_header(slide)
    draw_focus_bar(slide)
    for i, (num, title) in enumerate(BLOCKS):
        bx = MARGIN + (i % 2) * (BLOCK_W + BLOCK_GAP)
        by = BODY_Y + (i // 2) * (BLOCK_H + BLOCK_GAP)
        draw_activity_block(slide, bx, by, num, title)
    draw_squad_panel(slide)
    add_footer(slide, footer)
    return slide


def clone_slide(prs, source, footer):
    """Exact structural copy: deep-copy every shape element into a fresh
    blank slide. Safe here because all content is native shapes (no rels)."""
    new = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in list(source.shapes):
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    # swap the footer (last shape added by build order is the footer textbox)
    for shp in new.shapes:
        if shp.has_text_frame and shp.text_frame.text.startswith("MASTER"):
            shp.text_frame.paragraphs[0].runs[0].text = footer
    return new


def add_annotations(slide):
    """Deletable dashed accent 'HOW TO' callouts placed in spacious zones."""
    notes = [
        # (x, y, w, h, text) — sited inside diagram/notes zones
        (0.55, 2.10, 2.55, 0.62,
         "HOW TO — Set timing + lead coach in each block header. "
         "Keep the objective to one line."),
        (5.90, 2.10, 2.55, 0.62,
         "HOW TO — Copy players, arrows and zones from Slide 4 "
         "(Component Library) onto the court."),
        (0.55, 5.20, 2.55, 0.62,
         "HOW TO — Use the notes strip for constraints, progressions "
         "and setup / equipment reminders."),
        (5.90, 5.20, 2.55, 0.62,
         "HOW TO — Duplicate Slide 2 each week. Delete these blue "
         "notes — they are single text boxes."),
        (10.93, 3.30, 2.15, 0.80,
         "HOW TO — Update counts, names and status here for fast late "
         "changes. Tab moves between table cells."),
    ]
    for x, y, w, h, text in notes:
        box = add_box(slide.shapes, x, y, w, h, fill=C_WHITE,
                      line=C_ACCENT, line_w=1.0, dash=MSO_LINE_DASH_STYLE.DASH)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = IN(0.05)
        tf.margin_right = IN(0.05)
        tf.margin_top = IN(0.03)
        tf.margin_bottom = IN(0.03)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(7.5)
        r.font.bold = False
        r.font.color.rgb = C_ACCENT


# ---------------------------------------------------------------- component library

def _marker(sh, x, y, fill, line, text, text_color, d=MARKER_D, dash=None):
    """Plain (non-directional) player circle. d defaults to court-scale size."""
    m = add_box(sh, x, y, d, d, fill=fill, line=line, line_w=1.0,
                shape=MSO_SHAPE.OVAL, dash=dash)
    tf = m.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(7)
    r.font.bold = True
    r.font.color.rgb = text_color
    return m


def _marker_text(shp, label, text_color, size=7, lift=0.0):
    """lift > 0 nudges the centred text upward (extra bottom margin) —
    used by the chord puck whose visual centre sits above its bbox centre."""
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = 0
    tf.margin_bottom = IN(lift)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(label)
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = True
    f.color.rgb = text_color
    return shp


def draw_d_puck(sh, x, y, fill, line_color, label, text_color, dash=None):
    """
    Body-shape marker: a circle with a flattened back — the flat edge IS the
    shoulder line, the curved edge is the chest. Default orientation faces UP.
    A single CHORD autoshape: no grouping, rotates cleanly in one click, and
    reads as open/closed body shape instantly at court scale.
    """
    d = MARKER_D
    shp = sh.add_shape(MSO_SHAPE.CHORD, IN(x), IN(y), IN(d), IN(d))
    # chord cut across the lower edge: keep the sweep from 150° -> 30°
    # (through north), leaving a flat shoulder line at the back/south
    set_adj_angles(shp, 150, 30)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1.0)
        if dash is not None:
            shp.line.dash_style = dash
    _no_shadow(shp)
    return _marker_text(shp, label, text_color, lift=0.03)


def draw_vision_marker(slide, x, y, fill, line_color, label, text_color,
                       dash=None):
    """
    Player circle + translucent vision cone (70° wedge, default facing UP),
    grouped as one object. Rotate the group to set body orientation — the
    cone shows facing AND field of view. Placed on a defender it doubles as
    a cover-shadow tool (the passing lanes the body blocks).
    x, y = top-left of the player circle. Cone radius ≈3 m at court scale.
    """
    d = MARKER_D
    cone_r = d * 1.75                  # ≈3 m at court scale
    ccx = x + d / 2.0
    ccy = y + d / 2.0

    pie = slide.shapes.add_shape(MSO_SHAPE.PIE,
                                 IN(ccx - cone_r), IN(ccy - cone_r),
                                 IN(2 * cone_r), IN(2 * cone_r))
    set_adj_angles(pie, 235, 305)      # 70° wedge centred on north (270°)
    pie.fill.solid()
    cone_color = fill if fill not in (C_WHITE, None) else (line_color or C_DARK)
    pie.fill.fore_color.rgb = cone_color
    set_fill_alpha(pie, 20)
    pie.line.fill.background()
    _no_shadow(pie)

    circ = add_box(slide.shapes, x, y, d, d, fill=fill, line=line_color,
                   line_w=1.0, shape=MSO_SHAPE.OVAL, dash=dash)
    _marker_text(circ, label, text_color)

    return group_elements(slide, [pie, circ],
                          ccx - cone_r, ccy - cone_r,
                          2 * cone_r, 2 * cone_r, f"vm_{label}")


def draw_scan_zone(slide, x, y, fill, line_color, label, text_color,
                   dash=None):
    """
    Player circle + scanning awareness arc (180° sweep, default facing UP),
    grouped as one object. Represents the full area a player actively sweeps
    before receiving — much larger than the vision cone (≈5.5 m vs ≈3 m).

    SCAN ZONE:   180° sweep, dashed outline, 10% opacity — awareness
    VISION CONE:  70° sweep, solid,          20% opacity — focus
    BODY PUCK:    chord shape, no extra group — shoulder line

    Rotate the group to reorient. Combine all three (draw_combined_awareness)
    to show the full scanning picture on a single player.
    """
    d = MARKER_D
    scan_r = d * 3.2                   # ≈5.5 m at court scale
    ccx = x + d / 2.0
    ccy = y + d / 2.0

    scan_color = fill if fill not in (C_WHITE, None) else (line_color or C_DARK)
    s_pie = slide.shapes.add_shape(MSO_SHAPE.PIE,
                                   IN(ccx - scan_r), IN(ccy - scan_r),
                                   IN(2 * scan_r), IN(2 * scan_r))
    # 180° sweep from west (180°) clockwise through north (270°) to east (360°)
    set_adj_angles(s_pie, 180, 360)
    s_pie.fill.solid()
    s_pie.fill.fore_color.rgb = scan_color
    set_fill_alpha(s_pie, 10)
    s_pie.line.color.rgb = scan_color
    s_pie.line.width = Pt(1.0)
    s_pie.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    _no_shadow(s_pie)

    circ = add_box(slide.shapes, x, y, d, d, fill=fill, line=line_color,
                   line_w=1.0, shape=MSO_SHAPE.OVAL, dash=dash)
    _marker_text(circ, label, text_color)

    return group_elements(slide, [s_pie, circ],
                          ccx - scan_r, ccy - scan_r,
                          2 * scan_r, 2 * scan_r, f"scan_{label}")


def draw_combined_awareness(slide, x, y, fill, line_color, label, text_color):
    """
    All three awareness layers in one grouped element — the full scanning
    picture for a single player:
      outer arc   scanning zone (180°, dashed, 10% opacity, ≈5.5 m)
      inner wedge vision cone   ( 70°, solid,  20% opacity, ≈3 m)
      centre      body-shape puck
    Rotate the whole group to reorient the player.
    """
    d = MARKER_D
    scan_r = d * 3.2
    cone_r = d * 1.75
    ccx = x + d / 2.0
    ccy = y + d / 2.0
    scan_color = fill if fill not in (C_WHITE, None) else (line_color or C_DARK)

    # layer 1: scan zone (bottom — largest, most transparent)
    s_pie = slide.shapes.add_shape(MSO_SHAPE.PIE,
                                   IN(ccx - scan_r), IN(ccy - scan_r),
                                   IN(2 * scan_r), IN(2 * scan_r))
    set_adj_angles(s_pie, 180, 360)
    s_pie.fill.solid(); s_pie.fill.fore_color.rgb = scan_color
    set_fill_alpha(s_pie, 10)
    s_pie.line.color.rgb = scan_color
    s_pie.line.width = Pt(1.0)
    s_pie.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    _no_shadow(s_pie)

    # layer 2: vision cone (mid — focused facing direction)
    c_pie = slide.shapes.add_shape(MSO_SHAPE.PIE,
                                   IN(ccx - cone_r), IN(ccy - cone_r),
                                   IN(2 * cone_r), IN(2 * cone_r))
    set_adj_angles(c_pie, 235, 305)
    c_pie.fill.solid(); c_pie.fill.fore_color.rgb = scan_color
    set_fill_alpha(c_pie, 22)
    c_pie.line.fill.background()
    _no_shadow(c_pie)

    # layer 3: body-shape puck (top — shoulder line visible)
    puck = slide.shapes.add_shape(MSO_SHAPE.CHORD, IN(x), IN(y), IN(d), IN(d))
    set_adj_angles(puck, 150, 30)
    puck.fill.solid(); puck.fill.fore_color.rgb = fill if fill is not None else scan_color
    if line_color:
        puck.line.color.rgb = line_color; puck.line.width = Pt(1.0)
    else:
        puck.line.fill.background()
    _no_shadow(puck)
    _marker_text(puck, label, text_color, lift=0.03)

    return group_elements(slide, [s_pie, c_pie, puck],
                          ccx - scan_r, ccy - scan_r,
                          2 * scan_r, 2 * scan_r, f"combined_{label}")


def draw_long_cone(slide, apex_x, apex_y, color, name):
    """
    Standalone LONG vision cone — no player attached. A single 70° PIE
    wedge, default facing UP with the apex at the player's position.
    ≈8 m at court scale (vs ≈3 m on the attached cone) — a true sight
    line across the court. Drop it under any existing marker, rotate to
    aim, drag a corner handle to lengthen or shorten.
    """
    r = MARKER_D * 4.66                # ≈8 m at court scale
    pie = slide.shapes.add_shape(MSO_SHAPE.PIE,
                                 IN(apex_x - r), IN(apex_y - r),
                                 IN(2 * r), IN(2 * r))
    set_adj_angles(pie, 235, 305)      # 70° wedge centred on north (270°)
    pie.fill.solid()
    pie.fill.fore_color.rgb = color
    set_fill_alpha(pie, 20)
    pie.line.color.rgb = color
    pie.line.width = Pt(0.75)
    _no_shadow(pie)
    pie.name = name
    return pie


def _lib_label(sh, x, y, w, text):
    add_text(sh, x, y, w, 0.18, [[(text, 6.5, False, C_MID)]],
             align=PP_ALIGN.CENTER, wrap=False)


def _section_head(sh, x, y, text):
    add_box(sh, x, y + 0.03, 0.04, 0.16, fill=C_ACCENT, line=None)
    add_text(sh, x + 0.08, y, 2.6, 0.22, [[(text, 9, True, C_DARK)]],
             wrap=False)


def build_component_library(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sh = slide.shapes

    add_text(sh, MARGIN, 0.05, 10.0, 0.50, [
        [("COMPONENT LIBRARY", 16, True, C_DARK)],
        [("All markers are at true court scale (≈1.7 m on a 40 m court). "
          "Copy and paste directly onto any diagram — they will be "
          "proportionally correct.", 8.5, False, C_MID)],
    ])
    add_box(sh, 0, HEADER_H, SLIDE_W, RULE_H, fill=C_ACCENT, line=None)

    # ── Row 1: plain markers | equipment | lines | zones ────────────────────
    r1y = 0.72

    # PLAIN PLAYER MARKERS (no direction indicator — for quick annotation)
    px, py = 0.40, r1y
    _section_head(sh, px, py, "PLAYER MARKERS  (plain)")
    plain = [
        (C_ACCENT, None,  "A", C_WHITE, None,                   "Attacker"),
        (C_WHITE,  C_MID, "D", C_MID,  None,                   "Defender"),
        (C_DARK,   None,  "G", C_WHITE, None,                   "GK"),
        (C_WHITE,  C_MID, "N", C_MID,  MSO_LINE_DASH_STYLE.DASH, "Neutral"),
    ]
    for i, (fill, line, t, tc, dash, lbl) in enumerate(plain):
        mx = px + 0.06 + i * 0.50
        _marker(sh, mx, py + 0.26, fill, line, t, tc, dash=dash)
        _lib_label(sh, mx - 0.04, py + 0.44, 0.44, lbl)

    # EQUIPMENT (court-scale)
    ex, ey = 2.55, r1y
    _section_head(sh, ex, ey, "EQUIPMENT")
    # Cone ≈1.1 m tall
    add_box(sh, ex + 0.08, ey + 0.28, 0.09, 0.08, fill=C_MID, line=None,
            shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
    _lib_label(sh, ex + 0.01, ey + 0.44, 0.40, "Cone")
    # Pole ≈0.5 m wide × 2.5 m tall
    add_box(sh, ex + 0.62, ey + 0.24, 0.04, 0.20, fill=C_MID, line=None)
    _lib_label(sh, ex + 0.47, ey + 0.44, 0.40, "Pole")
    # Ball ≈0.85 m
    add_box(sh, ex + 1.10, ey + 0.30, 0.07, 0.07, fill=C_DARK, line=None,
            shape=MSO_SHAPE.OVAL)
    _lib_label(sh, ex + 0.95, ey + 0.44, 0.40, "Ball")

    # LINES
    lx, ly = 3.75, r1y
    _section_head(sh, lx, ly, "LINES")
    add_arrowhead(add_line(sh, lx + 0.05, ly + 0.38,
                           lx + 0.80, ly + 0.38, C_DARK, 1.5))
    _lib_label(sh, lx, ly + 0.44, 0.90, "Run")
    add_arrowhead(add_line(sh, lx + 1.00, ly + 0.38,
                           lx + 1.75, ly + 0.38,
                           C_DARK, 1.5, dash=MSO_LINE_DASH_STYLE.DASH))
    _lib_label(sh, lx + 0.95, ly + 0.44, 0.90, "Pass")
    add_arrowhead(add_line(sh, lx + 1.95, ly + 0.38,
                           lx + 2.70, ly + 0.38,
                           C_DARK, 1.5, dash=MSO_LINE_DASH_STYLE.ROUND_DOT))
    _lib_label(sh, lx + 1.90, ly + 0.44, 0.90, "Dribble")

    # ZONES & SCREENS
    zx, zy = 7.05, r1y
    _section_head(sh, zx, zy, "ZONES / SCREENS")
    zone = add_box(sh, zx + 0.05, zy + 0.26, 0.70, 0.24, fill=C_ACCENT, line=None)
    set_fill_alpha(zone, 18)
    _lib_label(sh, zx + 0.02, zy + 0.44, 0.76, "Zone")
    add_box(sh, zx + 1.00, zy + 0.34, 0.36, 0.06, fill=C_DARK, line=None,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _lib_label(sh, zx + 0.88, zy + 0.44, 0.62, "Screen")

    add_line(sh, MARGIN, 1.65, SLIDE_W - MARGIN, 1.65, C_LINE, 0.5)

    # ── Row 2: body shape & orientation ─────────────────────────────────────
    r2y = 1.72
    _section_head(sh, 0.40, r2y, "BODY SHAPE  —  shoulder-line pucks")

    puck_items = [
        (C_ACCENT, None,  "A", C_WHITE, None,                   "Attacker"),
        (C_WHITE,  C_MID, "D", C_MID,  None,                   "Defender"),
        (C_DARK,   None,  "G", C_WHITE, None,                   "GK"),
        (C_WHITE,  C_MID, "N", C_MID,  MSO_LINE_DASH_STYLE.DASH, "Neutral"),
    ]
    for i, (fill, line, t, tc, dash, lbl) in enumerate(puck_items):
        mx = 0.47 + i * 0.50
        draw_d_puck(sh, mx, r2y + 0.28, fill, line, t, tc, dash=dash)
        _lib_label(sh, mx - 0.06, r2y + 0.46, 0.44, lbl)

    add_text(sh, 0.40, r2y + 0.66, 2.55, 0.36, [[
        ("Flat edge = shoulders, curve = chest. One shape — rotate it to "
         "show open / closed body shape.", 7.5, False, C_MID),
    ]])

    _section_head(sh, 3.30, r2y, "VISION CONE  —  facing + field of view")
    cone_items = [
        (C_ACCENT, None,  "A", C_WHITE, "Attacker"),
        (C_WHITE,  C_MID, "D", C_MID,  "Defender"),
        (C_DARK,   None,  "G", C_WHITE, "GK"),
    ]
    for i, (fill, line, t, tc, lbl) in enumerate(cone_items):
        mx = 3.55 + i * 0.62
        draw_vision_marker(slide, mx, r2y + 0.42, fill, line, t, tc)
        _lib_label(sh, mx - 0.15, r2y + 0.60, 0.44, lbl)

    add_text(sh, 3.30, r2y + 0.80, 2.85, 0.30, [[
        ("Rotate the group: cone shows where the player faces and what "
         "they see. On a defender it doubles as a cover shadow.",
         7.5, False, C_MID),
    ]])

    # rotation instruction callout
    instr = add_box(sh, 6.55, r2y + 0.24, 3.30, 0.62,
                    fill=C_FAINT, line=C_LINE, line_w=0.75)
    tf = instr.text_frame
    tf.word_wrap = True
    tf.margin_left = IN(0.06); tf.margin_right = IN(0.06)
    tf.margin_top = IN(0.04); tf.margin_bottom = IN(0.04)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("ROTATION  ·  All orientation markers face UP by default. "
              "Select → drag the rotation handle, or Format Shape → "
              "Rotation and type the angle (0 = up, 90 = right, "
              "180 = down, 270 = left).")
    r.font.name = FONT; r.font.size = Pt(7.5); r.font.color.rgb = C_MID

    add_text(sh, 6.55, r2y + 0.90, 3.30, 0.30, [[
        ("LONG CONES (right) are separate — no player attached. Drop one "
         "under any marker, rotate to aim, drag a corner handle to "
         "lengthen.", 7.5, False, C_MID),
    ]])

    # LONG VISION CONES — standalone sight lines, drawn before the header
    # so the header text renders on top of the translucent wedges
    long_items = [(C_ACCENT, "Attacker"), (C_MID, "Defender"), (C_DARK, "GK")]
    for i, (color, lbl) in enumerate(long_items):
        ax = 10.55 + i * 1.05
        draw_long_cone(slide, ax, r2y + 0.95, color, f"Long cone {lbl}")
        _lib_label(sh, ax - 0.30, r2y + 0.99, 0.60, lbl)
    _section_head(sh, 10.00, r2y, "LONG VISION CONES  (≈8 m, standalone)")

    add_line(sh, MARGIN, 2.92, SLIDE_W - MARGIN, 2.92, C_LINE, 0.5)

    # ── Row 3: scanning awareness zone ──────────────────────────────────────
    # Shapes drawn FIRST so header text renders on top (z-order)
    r3y = 2.97
    player_sy = r3y + 0.35             # player top; scan group top ≈ 2.95 ✓

    scan_items = [
        (C_ACCENT, None,  "A", C_WHITE, "Attacker"),
        (C_WHITE,  C_MID, "D", C_MID,  "Defender"),
        (C_DARK,   None,  "G", C_WHITE, "GK"),
    ]
    _scan_r = MARKER_D * 3.2           # 0.448"
    # standalone examples
    for i, (fill, line, t, tc, lbl) in enumerate(scan_items):
        px_ = 0.65 + i * 1.10
        draw_scan_zone(slide, px_, player_sy, fill, line, t, tc)
        _lib_label(sh, px_ - 0.08, player_sy + MARKER_D + 0.06, 0.50, lbl)

    # combined: scan zone + vision cone + body puck in one group
    comb_x = 4.08
    draw_combined_awareness(slide, comb_x, player_sy, C_ACCENT, None, "A", C_WHITE)
    _lib_label(sh, comb_x - 0.14, player_sy + MARKER_D + 0.06, 0.58, "Combined")

    # section header drawn on top of the translucent zones
    _section_head(sh, 0.40, r3y,
                  "SCANNING AWARENESS ZONE  (180°, ≈5.5 m at court scale)")

    # explanation (right of the examples)
    add_text(sh, 5.05, r3y + 0.05, 7.90, 1.00, [
        [("SCANNING AWARENESS  ·  VISION CONE  ·  BODY SHAPE", 8.5, True, C_DARK)],
        [("Scan zone (180°, dashed, very light) = the full area a player "
          "actively sweeps before receiving — peripheral awareness, "
          "not just facing direction.", 8, False, C_MID)],
        [("", 3, False, C_DARK)],
        [("Vision cone (70°, solid) = immediate focus and field of view. "
          "Scan zone + vision cone together show the difference between what "
          "a player checks and what they commit to.", 8, False, C_MID)],
        [("", 3, False, C_DARK)],
        [("COMBINED example (left): all three layers in one group. "
          "Rotate the whole group to orient the player. "
          "Use standalone or combined — each is a separate object.",
          8, False, C_DARK)],
    ])

    add_line(sh, MARGIN, 4.15, SLIDE_W - MARGIN, 4.15, C_LINE, 0.5)

    # ── Row 4: Team A + Team B numbered sets ────────────────────────────────
    r4y = 4.21

    add_text(sh, 0.40, r4y, 2.4, 0.20, [[
        ("TEAM A  ", 8, True, C_ACCENT),
        ("your team", 7.5, False, C_MID),
    ]])
    for n in range(1, 6):
        draw_d_puck(sh, 0.45 + (n - 1) * 0.38, r4y + 0.24,
                    C_ACCENT, None, str(n), C_WHITE)

    add_text(sh, 2.75, r4y, 3.0, 0.20, [[
        ("TEAM B  ", 8, True, C_DARK),
        ("opposition — change fill colour to match", 7.5, False, C_MID),
    ]])
    for n in range(1, 6):
        draw_d_puck(sh, 2.80 + (n - 1) * 0.38, r4y + 0.24,
                    C_DARK, None, str(n), C_WHITE)

    add_text(sh, 5.50, r4y, 7.50, 0.55, [
        [("SCALE + COLOUR", 7, True, C_MID)],
        [("Markers are at true court scale. Team pucks carry body shape: "
          "rotate any puck to angle the shoulders. "
          "To recolour Team B: select all 5 → Format Shape → Fill.",
          8, False, C_DARK)],
    ])

    add_line(sh, MARGIN, 4.85, SLIDE_W - MARGIN, 4.85, C_LINE, 0.5)

    # ── Row 5: courts ───────────────────────────────────────────────────────
    cy0 = 4.92
    _section_head(sh, 0.40, cy0, "FULL COURT  (40m × 20m)")
    draw_futsal_court(slide, 0.55, cy0 + 0.28, 3.80, "Library full court")
    _section_head(sh, 5.60, cy0, "HALF COURT  (20m × 20m)")
    draw_futsal_half_court(slide, 5.75, cy0 + 0.28, 1.90,
                           "Library half court")

    add_text(sh, 8.10, cy0 + 0.28, 4.85, 2.00, [
        [("USAGE", 7, True, C_MID)],
        [("Courts are grouped — one click selects the whole court. "
          "Copy and paste into any diagram zone, resize to fit. "
          "Ungroup (Ctrl+Shift+G) to edit individual markings.", 8, False, C_DARK)],
        [("", 3, False, C_DARK)],
        [("All shapes are unlocked. Numbers inside markers are editable — "
          "double-click to change. To recolour: select shapes → "
          "Format Shape → Fill.", 8, False, C_DARK)],
    ])

    add_footer(slide, "COMPONENT LIBRARY  —  court-scale tactical symbols  ·  "
                      "scan zone (180°, 5.5 m) · vision cone (70°, 3 m) · "
                      "long cone (8 m, standalone) · body puck · "
                      "rotate to orient · combined group available")
    return slide


# ---------------------------------------------------------------- session control board
#
# No session plan in any sport shows TIME, LOAD, BALANCE and the REVIEW
# LOOP on one surface. This board does. It treats the session as an
# instrument, not a document:
#
#   1. TEMPO & LOAD TIMELINE  — width = minutes; a draggable intensity
#      curve built from diamond nodes wired to live connectors (drag a
#      node, the curve follows — native PowerPoint behaviour).
#   2. GAME-MOMENT FINGERPRINT — futsal's four moments + set plays vs
#      each activity. The session's tactical balance, visible at arm's
#      length before a single cone goes down.
#   3. ROTATION BOARD          — who is in which team per block.
#   4. REVIEW LOOP + AUDIT     — closed before leaving the hall; the
#      carry-forward seeds next week's session theme.

def _slide_masthead(slide, title, subtitle):
    add_text(slide.shapes, MARGIN, 0.05, 12.9, 0.50, [
        [(title, 16, True, C_DARK)],
        [(subtitle, 8.5, False, C_MID)],
    ])
    add_box(slide.shapes, 0, HEADER_H, SLIDE_W, RULE_H,
            fill=C_ACCENT, line=None)


def _field_box(sh, x, y, w, h, label, value="[   ]"):
    add_box(sh, x, y, w, h, fill=C_WHITE, line=C_LINE, line_w=0.75)
    add_text(sh, x + 0.04, y + 0.02, w - 0.08, 0.16,
             [[(label, 6.5, True, C_MID)]], wrap=False)
    add_text(sh, x + 0.04, y + h - 0.30, w - 0.08, 0.26,
             [[(value, 11, True, C_DARK)]])


def build_control_board(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sh = slide.shapes

    _slide_masthead(
        slide, "SESSION CONTROL BOARD",
        "The session as an instrument, not a document — time, load, "
        "tactical balance, rotations and the review loop on one surface. "
        "Plan it here, audit it here, seed next week here.")

    # ── Band A: tempo & load timeline ───────────────────────────────────────
    _section_head(sh, 0.40, 0.68, "TEMPO + LOAD TIMELINE  (width = time)")
    add_text(sh, 4.10, 0.66, 8.95, 0.40, [[
        ("Bars are minutes — drag a bar edge to your real timings. "
         "Diamonds = planned intensity: drag one and the curve stays "
         "attached. After the session, copy the diamonds, recolour grey, "
         "set actual — instant planned vs actual.", 7.5, False, C_MID),
    ]])

    PLOT_X, PLOT_R, TOTAL_MIN = 1.10, 12.95, 90
    PLOT_W = PLOT_R - PLOT_X

    def TX(m):                       # session minute -> x
        return PLOT_X + m * (PLOT_W / TOTAL_MIN)

    def TY(rpe):                     # RPE 1-5 -> y
        return 2.12 - (rpe - 1) * 0.25

    # intensity grid
    for r in range(1, 6):
        add_line(sh, PLOT_X, TY(r), PLOT_R, TY(r), C_LINE, 0.5)
        add_text(sh, 0.55, TY(r) - 0.08, 0.50, 0.16,
                 [[(f"RPE {r}", 6, False, C_MID)]],
                 align=PP_ALIGN.RIGHT, wrap=False)

    # activity duration bars (proportional defaults — coach drags edges)
    bar_y, bar_h = 2.30, 0.28
    blocks_t = [(0, 15, "01  WARM UP · 15'"),
                (20, 20, "02  TECHNICAL · 20'"),
                (45, 20, "03  TACTICAL · 20'"),
                (70, 20, "04  GAME · 20'")]
    for start, mins, lbl in blocks_t:
        bar = add_box(sh, TX(start), bar_y, TX(start + mins) - TX(start),
                      bar_h, fill=C_ACCENT, line=C_ACCENT, line_w=0.75)
        set_fill_alpha(bar, 15)
        add_text(sh, TX(start) + 0.02, bar_y + 0.055,
                 TX(start + mins) - TX(start) - 0.04, 0.18,
                 [[(lbl, 7, True, C_DARK)]], wrap=False)
    for start, end in ((15, 20), (40, 45), (65, 70)):     # rest gaps
        add_box(sh, TX(start), bar_y, TX(end) - TX(start), bar_h,
                fill=C_LITE, line=None)
        add_text(sh, TX(start), bar_y + 0.055, TX(end) - TX(start), 0.18,
                 [[("R", 6.5, True, C_MID)]], align=PP_ALIGN.CENTER,
                 wrap=False)

    # minute axis: 5-min minor ticks, 15-min labelled
    ax_y = bar_y + bar_h
    for m in range(0, TOTAL_MIN + 1, 5):
        major = m % 15 == 0
        add_line(sh, TX(m), ax_y, TX(m), ax_y + (0.05 if major else 0.03),
                 C_MID if major else C_LINE, 0.5)
        if major:
            add_text(sh, TX(m) - 0.20, ax_y + 0.05, 0.40, 0.14,
                     [[(str(m), 6, False, C_MID)]],
                     align=PP_ALIGN.CENTER, wrap=False)

    # intensity curve: diamond nodes + connectors that stay attached
    nodes_def = [(0, 2), (15, 3), (20, 2.5), (40, 4),
                 (45, 3.5), (65, 4.5), (70, 4), (90, 5)]
    nd = 0.12
    nodes = []
    for m, rpe in nodes_def:
        node = add_box(sh, TX(m) - nd / 2, TY(rpe) - nd / 2, nd, nd,
                       fill=C_ACCENT, line=None, shape=MSO_SHAPE.DIAMOND)
        nodes.append(node)
    for a, b in zip(nodes, nodes[1:]):
        seg = add_line(sh, 0, 0, 1, 1, C_ACCENT, 1.5)
        seg.begin_connect(a, 3)      # right vertex of left node
        seg.end_connect(b, 1)        # left vertex of right node

    add_line(sh, MARGIN, 2.88, SLIDE_W - MARGIN, 2.88, C_LINE, 0.5)

    # ── Band B left: game-moment fingerprint ────────────────────────────────
    _section_head(sh, 0.40, 2.98, "GAME-MOMENT FINGERPRINT")
    add_text(sh, 0.42, 3.18, 4.40, 0.16, [[
        ("Where does this session actually live? Fill before you coach.",
         7, False, C_MID)]], wrap=False)

    moments = ("ATT ORG", "DEF ORG", "ATT→DEF", "DEF→ATT", "SET PLAY")
    col_x0, col_w = 1.62, 0.64
    for c, mom in enumerate(moments):
        add_text(sh, col_x0 + c * col_w, 3.40, col_w, 0.14,
                 [[(mom, 6, True, C_MID)]], align=PP_ALIGN.CENTER,
                 wrap=False)
    fp_rows = ("01  WARM UP", "02  TECHNICAL", "03  TACTICAL", "04  GAME")
    sq = 0.28
    for r_i, lbl in enumerate(fp_rows):
        ry = 3.60 + r_i * 0.38
        add_text(sh, 0.42, ry + 0.04, 1.18, 0.18,
                 [[(lbl, 7.5, True, C_DARK)]], wrap=False)
        for c in range(5):
            add_box(sh, col_x0 + c * col_w + (col_w - sq) / 2, ry,
                    sq, sq, fill=C_WHITE, line=C_LINE, line_w=0.75)
    tot_y = 3.60 + 4 * 0.38 + 0.04
    add_text(sh, 0.42, tot_y + 0.02, 1.18, 0.18,
             [[("MINUTES", 7, True, C_MID)]], wrap=False)
    for c in range(5):
        add_text(sh, col_x0 + c * col_w, tot_y, col_w, 0.20,
                 [[("[  ]", 8, True, C_DARK)]], align=PP_ALIGN.CENTER,
                 wrap=False)
    leg_y = tot_y + 0.30
    lsq = add_box(sh, 0.42, leg_y, 0.14, 0.14, fill=C_ACCENT, line=None)
    add_text(sh, 0.60, leg_y - 0.02, 1.10, 0.18,
             [[("PRIMARY", 6.5, False, C_MID)]], wrap=False)
    add_box(sh, 1.40, leg_y, 0.14, 0.14, fill=C_LITE, line=C_LINE,
            line_w=0.75)
    add_text(sh, 1.58, leg_y - 0.02, 3.30, 0.18,
             [[("SECONDARY  ·  click a square → Fill to mark",
                6.5, False, C_MID)]], wrap=False)

    # ── Band B middle: rotation board ───────────────────────────────────────
    _section_head(sh, 5.15, 2.98, "ROTATION BOARD")
    add_text(sh, 5.17, 3.18, 3.30, 0.16, [[
        ("A / B = team   N = neutral   G = GK   R = rest",
         7, False, C_MID)]], wrap=False)

    positions = (["GK"] * 2 + ["FIXO"] * 3 + ["ALA"] * 6 + ["PIVOT"] * 3)
    rrows = 1 + len(positions)
    rt_y, rt_h = 3.44, 3.30
    gframe = sh.add_table(rrows, 6, IN(5.15), IN(rt_y), IN(3.30), IN(rt_h))
    tbl = gframe.table
    tbl.first_row = False
    tbl.horz_banding = False
    tblPr = tbl._tbl.tblPr
    for el in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(el)
    tsid = tblPr.makeelement(qn("a:tableStyleId"), {})
    tsid.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
    tblPr.append(tsid)
    for c, cw in enumerate((0.42, 1.04, 0.46, 0.46, 0.46, 0.46)):
        tbl.columns[c].width = IN(cw)
    row_h = IN(rt_h / rrows)
    for r in tbl.rows:
        r.height = row_h
    for c, head in enumerate(("POS", "PLAYER", "01", "02", "03", "04")):
        _style_cell(tbl.cell(0, c), head, 6.5, True, C_DARK, C_LITE,
                    align=PP_ALIGN.CENTER if c >= 2 else PP_ALIGN.LEFT)
    for i, pos in enumerate(positions, start=1):
        tint = C_FAINT if pos in ("FIXO", "PIVOT") else C_WHITE
        _style_cell(tbl.cell(i, 0), pos, 6.5, True, C_DARK, tint)
        _style_cell(tbl.cell(i, 1), "[Player]", 6.5, False, C_MID, tint)
        for c in range(2, 6):
            _style_cell(tbl.cell(i, c), "", 6.5, True, C_DARK, tint,
                        align=PP_ALIGN.CENTER)
    add_text(sh, 5.15, 7.06, 3.30, 0.20, [[
        ("Empty cells are visible before the session — nobody hides.",
         6.5, False, C_MID)]], wrap=False)

    # ── Band B right: review loop ───────────────────────────────────────────
    _section_head(sh, 8.75, 2.98, "REVIEW LOOP  —  close it in the hall")
    loop_items = [
        ("WHAT WORKED", "[Insert — be specific]", C_LINE),
        ("WHAT TO CHANGE", "[Insert — one constraint to adjust]", C_LINE),
        ("CARRY FORWARD → NEXT SESSION", "[Insert — paste into next "
         "week's SESSION THEME]", C_ACCENT),
    ]
    ly = 3.28
    for label, value, edge in loop_items:
        add_box(sh, 8.75, ly, 4.20, 0.60, fill=C_WHITE, line=edge,
                line_w=1.0)
        add_text(sh, 8.81, ly + 0.04, 4.08, 0.50, [
            [(label, 7, True, C_ACCENT if edge is C_ACCENT else C_MID)],
            [(value, 8, False, C_DARK)],
        ])
        if edge is not C_ACCENT:
            add_arrowhead(add_line(sh, 10.85, ly + 0.60, 10.85, ly + 0.72,
                                   C_MID, 1.0))
        ly += 0.72

    # ── Band C left: session audit ──────────────────────────────────────────
    _section_head(sh, 0.40, 5.62, "SESSION AUDIT  —  the numbers nobody tracks")
    audit = (("BALL-ROLLING TIME  (%)", 0.42, 5.92),
             ("COACH TALK TIME  (min)", 2.68, 5.92),
             ("PLANNED AVG INTENSITY", 0.42, 6.56),
             ("ACTUAL AVG INTENSITY", 2.68, 6.56))
    for label, fx, fy in audit:
        _field_box(sh, fx, fy, 2.16, 0.56, label)

    # ── Band C right: next session seed ─────────────────────────────────────
    _section_head(sh, 8.75, 5.62, "NEXT SESSION SEED  —  first three decisions")
    for i in range(3):
        sy = 5.92 + i * 0.42
        add_text(sh, 8.78, sy, 4.17, 0.30, [[
            (f"{i + 1}   ", 10, True, C_ACCENT),
            ("[Insert decision]", 8.5, False, C_DARK),
        ]])
        add_line(sh, 8.78, sy + 0.32, 12.95, sy + 0.32, C_LINE, 0.5)

    add_footer(slide, "SESSION CONTROL BOARD  —  width = time · drag "
                      "intensity diamonds, the curve follows · fingerprint "
                      "shows tactical balance · review loop seeds next week")
    return slide


# ---------------------------------------------------------------- on-court cards

def build_oncourt_cards(prs):
    """Print → cut along the dashed lines → pocket. Four cards, one per
    block, in session order — type sized to read at arm's length in a
    loud hall, mid-session, without a clipboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sh = slide.shapes

    _slide_masthead(
        slide, "ON-COURT CARDS",
        "Print this slide, cut along the dashed lines: four pocket cards, "
        "one per block, in session order. Three cues maximum per card — "
        "if it needs more words, it belongs on the plan, not on court.")

    cards = ((0.30, 0.85, "01", "WARM UP"),
             (7.00, 0.85, "02", "TECHNICAL / POSSESSION"),
             (0.30, 4.15, "03", "TACTICAL ACTIVITY"),
             (7.00, 4.15, "04", "GAME"))
    cw, ch = 6.10, 3.05

    for cx, cy, num, title in cards:
        add_box(sh, cx, cy, cw, ch, fill=C_WHITE, line=C_LINE, line_w=1.0)
        add_box(sh, cx, cy, 0.07, ch, fill=C_ACCENT, line=None)
        add_text(sh, cx + 0.15, cy + 0.06, 0.75, 0.50,
                 [[(num, 26, True, C_ACCENT)]], wrap=False)
        add_text(sh, cx + 0.90, cy + 0.20, 3.55, 0.32,
                 [[(title, 14, True, C_DARK)]], wrap=False)
        chip = add_box(sh, cx + 4.60, cy + 0.16, 1.32, 0.38,
                       fill=C_WHITE, line=C_ACCENT, line_w=1.25,
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _marker_text(chip, "[Insert mins]", C_ACCENT, size=10)
        add_text(sh, cx + 0.20, cy + 0.66, 1.20, 0.18,
                 [[("CUES", 7.5, True, C_MID)]], wrap=False)
        for i in range(3):
            qy = cy + 0.88 + i * 0.40
            add_text(sh, cx + 0.20, qy, 3.60, 0.34, [[
                (f"{i + 1}  ", 13, True, C_ACCENT),
                ("[Insert cue — five words max]", 12, True, C_DARK),
            ]], wrap=False)
        draw_futsal_court(slide, cx + 3.95, cy + 0.95, 1.90,
                          f"Card court {num}")
        add_line(sh, cx + 0.20, cy + 2.32, cx + cw - 0.15, cy + 2.32,
                 C_LINE, 0.5)
        add_text(sh, cx + 0.20, cy + 2.40, 3.90, 0.55, [[
            ("WATCH FOR   ", 7.5, True, C_MID),
            ("[Insert the trigger that makes you stop play]",
             9.5, False, C_DARK),
        ]])
        add_text(sh, cx + 4.30, cy + 2.40, 1.65, 0.30, [[
            ("WORK : REST   ", 7.5, True, C_MID),
            ("[   ]", 9.5, True, C_DARK),
        ]], align=PP_ALIGN.RIGHT, wrap=False)

    # cut guides
    add_line(sh, 6.70, 0.78, 6.70, 7.28, C_MID, 0.75,
             dash=MSO_LINE_DASH_STYLE.DASH)
    add_line(sh, 0.20, 4.025, 13.13, 4.025, C_MID, 0.75,
             dash=MSO_LINE_DASH_STYLE.DASH)
    add_text(sh, 6.46, 0.60, 0.50, 0.16, [[("CUT", 6, True, C_MID)]],
             align=PP_ALIGN.CENTER, wrap=False)
    add_text(sh, 13.00, 3.93, 0.34, 0.16, [[("CUT", 6, True, C_MID)]],
             wrap=False)

    add_footer(slide, "ON-COURT CARDS  —  print · cut · pocket  ·  one "
                      "card per block  ·  three cues max  ·  readable at "
                      "arm's length")
    return slide


# ---------------------------------------------------------------- main

def main():
    prs = Presentation()
    prs.slide_width = IN(SLIDE_W)
    prs.slide_height = IN(SLIDE_H)

    master = build_session_slide(
        prs, "MASTER TEMPLATE  —  keep clean  ·  duplicate Slide 2 for "
             "each session")
    clone_slide(prs, master,
                "WORKING COPY  —  ready to edit  ·  duplicate weekly")
    guide = clone_slide(prs, master,
                        "GUIDE  —  annotated example  ·  delete the blue "
                        "notes after reading")
    add_annotations(guide)
    build_component_library(prs)
    build_control_board(prs)
    build_oncourt_cards(prs)

    out = "futsal_session_template.pptx"
    prs.save(out)
    print(f"Saved {out}: {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
